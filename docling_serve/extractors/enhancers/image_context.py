"""Image-context: describe document images with a vision agent.

Two consumers share one core:

* :class:`ImageContextEnhancer` — deep-extraction pass that walks the deep
  document's image assets, describes each with the Bedrock vision model, and
  writes the context back onto assets/elements plus an ``image-context.json``
  sidecar.
* :func:`describe_file_images` — stateless entry point (used by the
  ``/v1/images/context`` endpoint) that pulls the significant raster images
  straight out of an uploaded file (PPTX picture shapes, PDF embedded images,
  raw rasters), de-duplicates them, and returns vision descriptions — so the
  standard ingest path can make slide-image content searchable/graphable
  without running a full deep-extraction bundle.
"""

from __future__ import annotations

import hashlib
import io
import logging
from pathlib import Path
from typing import Any

from docling_serve.deep_document.artifact_writer import write_json
from docling_serve.extractors.base import ExtractionContext, ExtractorResult
from docling_serve.extractors.enhancers.base import EnhancementResult, Enhancer
from docling_serve.providers import (
    BedrockUnavailableError,
    VisionMessage,
    get_bedrock_provider,
)
from docling_serve.settings import docling_serve_settings

_log = logging.getLogger(__name__)

_RASTER_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
# Skip decorative assets (logos, bullets, divider art): tiny files carry no
# document knowledge but each costs a vision call.
_MIN_IMAGE_BYTES = 8_192
_MIN_IMAGE_DIMENSION = 64

DEFAULT_PROMPT = (
    "You are analyzing an image extracted from a document. Describe what the "
    "image shows in 2-4 sentences so a reader who cannot see it understands its "
    "content and purpose. Transcribe any visible text, labels, or numbers. If it "
    "is a chart, diagram, schematic, table, or form, say so and summarize the key "
    "information. Be specific and factual; do not speculate."
)


class ImageContextEnhancer(Enhancer):
    name = "image_context"
    aliases = ("images", "image-context", "image_captions", "vision_context")

    def applies(self, ctx: ExtractionContext, document: dict[str, Any]) -> bool:
        provider = get_bedrock_provider()
        if not provider.enabled:
            return False
        return bool(_image_assets(document))

    def enhance(
        self,
        ctx: ExtractionContext,
        document: dict[str, Any],
        *,
        base_result: ExtractorResult,
    ) -> EnhancementResult:
        provider = get_bedrock_provider()
        result = EnhancementResult(name=self.name)
        if not provider.enabled:
            result.notes.append("bedrock_disabled")
            return result

        max_images = int(getattr(docling_serve_settings, "enhancement_max_images", 40))
        assets = _image_assets(document)[:max_images]

        context_by_asset: dict[str, str] = {}
        records: list[dict[str, Any]] = []
        for asset in assets:
            rel_path = str(asset.get("path") or "")
            image_path = ctx.bundle_dir / rel_path
            if not image_path.is_file():
                continue
            try:
                image_bytes = image_path.read_bytes()
                text = provider.converse(
                    messages=[
                        VisionMessage(
                            text=DEFAULT_PROMPT,
                            images=[image_bytes],
                            image_format=_image_format(image_path),
                        )
                    ]
                )
            except BedrockUnavailableError as err:
                _log.warning("Image context failed for %s: %s", rel_path, err)
                result.notes.append(f"{rel_path}: {err}")
                continue
            asset_id = str(asset.get("assetId") or rel_path)
            asset["context"] = text
            context_by_asset[asset_id] = text
            records.append({"assetId": asset_id, "path": rel_path, "context": text})

        if not records:
            result.notes.append("no_images_described")
            return result

        _annotate_elements(document, context_by_asset)

        sidecar = ctx.bundle_dir / "image-context.json"
        write_json(
            sidecar,
            {
                "model": {"provider": "bedrock", "modelId": provider.vision_model},
                "count": len(records),
                "images": records,
            },
        )

        result.applied = True
        result.artifacts.append(sidecar.relative_to(ctx.bundle_dir).as_posix())
        result.manifest_extra = {
            "imageContext": {
                "count": len(records),
                "sidecar": sidecar.relative_to(ctx.bundle_dir).as_posix(),
                "model": {"provider": "bedrock", "modelId": provider.vision_model},
            }
        }
        return result


def _image_assets(document: dict[str, Any]) -> list[dict[str, Any]]:
    assets = document.get("assets")
    if not isinstance(assets, list):
        return []
    images: list[dict[str, Any]] = []
    for asset in assets:
        if not isinstance(asset, dict):
            continue
        content_type = str(asset.get("contentType") or "")
        path = str(asset.get("path") or "")
        if content_type.startswith("image/") or Path(path).suffix.lower() in _RASTER_SUFFIXES:
            images.append(asset)
    return images


def _annotate_elements(document: dict[str, Any], context_by_asset: dict[str, str]) -> None:
    units = (document.get("document") or {}).get("units") or []
    for unit in units:
        if not isinstance(unit, dict):
            continue
        for element in (unit.get("content") or {}).get("elements") or []:
            if not isinstance(element, dict):
                continue
            asset_ref = element.get("assetRef")
            if asset_ref and asset_ref in context_by_asset:
                element["context"] = context_by_asset[asset_ref]


def _image_format(path: Path) -> str:
    suffix = path.suffix.lower().lstrip(".")
    return "jpeg" if suffix in {"jpg", "jpeg"} else suffix or "png"


# --------------------------------------------------------------------------- #
# Stateless core (shared by the /v1/images/context endpoint)                  #
# --------------------------------------------------------------------------- #


class ImageContextUnavailable(RuntimeError):
    """Raised when image description cannot run (Bedrock disabled, no images)."""


def extract_file_images(
    path: Path, *, max_images: int | None = None
) -> list[dict[str, Any]]:
    """Pull significant raster images out of a document file.

    Supports PPTX (picture shapes, slide-attributed), PDF (embedded image
    objects, page-attributed), and raw raster files. Decorative images are
    filtered by size, and repeated images (template logos/banners) are
    de-duplicated by content hash — on real decks this collapses most of the
    image count. Returns ``[{label, page, format, bytes}]`` in document order.
    """
    limit = max_images or int(getattr(docling_serve_settings, "enhancement_max_images", 40))
    suffix = path.suffix.lower()
    if suffix in {".pptx", ".ppt"}:
        raw = _pptx_images(path)
    elif suffix == ".pdf":
        raw = _pdf_images(path)
    elif suffix in _RASTER_SUFFIXES:
        raw = [{"label": path.name, "page": 1, "format": _image_format(path), "bytes": path.read_bytes()}]
    else:
        raw = []

    seen: set[str] = set()
    out: list[dict[str, Any]] = []
    for item in raw:
        data = item["bytes"]
        if len(data) < _MIN_IMAGE_BYTES:
            continue
        digest = hashlib.sha256(data).hexdigest()
        if digest in seen:
            continue
        seen.add(digest)
        if not _meets_min_dimensions(data):
            continue
        out.append(item)
        if len(out) >= limit:
            break
    return out


def describe_file_images(
    path: Path, *, max_images: int | None = None
) -> list[dict[str, Any]]:
    """Extract + vision-describe a file's images; returns ``[{label, page, description}]``.

    Raises :class:`ImageContextUnavailable` when Bedrock is disabled so callers
    can degrade to an empty response instead of failing the request.
    """
    provider = get_bedrock_provider()
    if not provider.enabled:
        raise ImageContextUnavailable("bedrock_disabled")
    images = extract_file_images(path, max_images=max_images)
    records: list[dict[str, Any]] = []
    for item in images:
        try:
            text = provider.converse(
                messages=[
                    VisionMessage(
                        text=DEFAULT_PROMPT,
                        images=[item["bytes"]],
                        image_format=item.get("format") or "png",
                    )
                ]
            )
        except BedrockUnavailableError as err:
            _log.warning("Image context failed for %s: %s", item.get("label"), err)
            continue
        records.append(
            {"label": item.get("label"), "page": item.get("page"), "description": text}
        )
    return records


def _pptx_images(path: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:  # pragma: no cover - python-pptx is a declared dep
        return []

    def walk(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)

    out: list[dict[str, Any]] = []
    try:
        prs = Presentation(str(path))
    except Exception as err:
        _log.warning("python-pptx could not open %s: %s", path, err)
        return []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in walk(slide.shapes):
            if shape.shape_type not in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                continue
            try:
                blob = shape.image.blob
                ext = (shape.image.ext or "png").lower()
            except Exception:  # unloadable media (e.g. WMF without a loader)
                continue
            if ext not in {"png", "jpg", "jpeg", "gif", "webp"}:
                converted = _convert_to_png(blob)
                if converted is None:
                    continue
                blob, ext = converted, "png"
            out.append(
                {
                    "label": f"slide {slide_no}: {getattr(shape, 'name', 'image')}",
                    "page": slide_no,
                    "format": "jpeg" if ext in {"jpg", "jpeg"} else ext,
                    "bytes": blob,
                }
            )
    return out


def _pdf_images(path: Path) -> list[dict[str, Any]]:
    try:
        import pypdfium2 as pdfium
        import pypdfium2.raw as pdfium_c
    except ImportError:  # pragma: no cover - pypdfium2 is a docling dep
        return []
    out: list[dict[str, Any]] = []
    try:
        pdf = pdfium.PdfDocument(str(path))
    except Exception as err:
        _log.warning("pypdfium2 could not open %s: %s", path, err)
        return []
    try:
        for page_index in range(len(pdf)):
            page = pdf[page_index]
            for obj_index, obj in enumerate(
                page.get_objects(filter=(pdfium_c.FPDF_PAGEOBJ_IMAGE,)), start=1
            ):
                try:
                    pil_image = obj.get_bitmap(render=False).to_pil()
                except Exception:
                    continue
                buffer = io.BytesIO()
                pil_image.convert("RGB").save(buffer, format="PNG", optimize=True)
                out.append(
                    {
                        "label": f"page {page_index + 1}: image {obj_index}",
                        "page": page_index + 1,
                        "format": "png",
                        "bytes": buffer.getvalue(),
                    }
                )
    finally:
        pdf.close()
    return out


def _convert_to_png(blob: bytes) -> bytes | None:
    try:
        from PIL import Image

        image = Image.open(io.BytesIO(blob))
        buffer = io.BytesIO()
        image.convert("RGB").save(buffer, format="PNG", optimize=True)
        return buffer.getvalue()
    except Exception:
        return None


def _meets_min_dimensions(data: bytes) -> bool:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
        return width >= _MIN_IMAGE_DIMENSION and height >= _MIN_IMAGE_DIMENSION
    except Exception:
        return True  # if unreadable here, let the vision call decide
