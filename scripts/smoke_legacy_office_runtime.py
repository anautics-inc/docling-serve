"""Production-image smoke test for real LibreOffice legacy round trips."""

from __future__ import annotations

import asyncio
import os
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path

# This is an air-gapped deployment smoke. Set offline policy before importing
# Docling/Hugging Face modules so their module-level constants honor it.
os.environ.setdefault("HF_HUB_OFFLINE", "1")
os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from docling.datamodel.base_models import DocumentStream
from docling.datamodel.service.chunking import HierarchicalChunkerOptions
from docling.datamodel.service.options import ConvertDocumentsOptions
from docling.datamodel.service.targets import InBodyTarget
from docling.datamodel.service.tasks import TaskType
from docling_jobkit.orchestrators.local.orchestrator import LocalOrchestratorConfig

from docling_serve.legacy_office import (
    LibreOfficeHeadlessConverter,
    build_converter_manager,
    check_legacy_office_capability,
)
from docling_serve.local_orchestrator import DoclingServeLocalOrchestrator
from docling_serve.orchestrator_factory import _build_cm_config


def _build_modern_fixtures(root: Path) -> dict[str, Path]:
    docx_path = root / "legacy-smoke.docx"
    document = Document()
    document.add_paragraph("legacy office smoke")
    document.save(docx_path)

    xlsx_path = root / "legacy-smoke.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "legacy office smoke"
    workbook.save(xlsx_path)

    pptx_path = root / "legacy-smoke.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "legacy office smoke"
    presentation.save(pptx_path)
    return {"doc": docx_path, "xls": xlsx_path, "ppt": pptx_path}


def main() -> None:
    executable = check_legacy_office_capability()
    subprocess.run(
        [str(executable), "--headless", "--version"],
        check=True,
        timeout=30,
    )
    with tempfile.TemporaryDirectory(prefix="legacy-office-smoke-") as raw_root:
        root = Path(raw_root)
        modern = _build_modern_fixtures(root)
        legacy_dir = root / "legacy"
        legacy_dir.mkdir()
        profile_dir = root / "profile"
        profile_dir.mkdir(mode=0o700)
        export_env = {
            "HOME": str(root),
            "TMPDIR": str(root),
            "XDG_CACHE_HOME": str(root / "cache"),
            "XDG_CONFIG_HOME": str(root / "config"),
            "LANG": "C.UTF-8",
            "LC_ALL": "C.UTF-8",
        }
        filters = {
            "doc": "doc:MS Word 97",
            "xls": "xls:MS Excel 97",
            "ppt": "ppt:MS PowerPoint 97",
        }
        for suffix, source in modern.items():
            subprocess.run(
                [
                    str(executable),
                    "--headless",
                    "--nologo",
                    "--nodefault",
                    "--nolockcheck",
                    "--nofirststartwizard",
                    f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                    "--convert-to",
                    filters[suffix],
                    "--outdir",
                    str(legacy_dir),
                    str(source),
                ],
                check=True,
                timeout=60,
                env={**os.environ, **export_env},
            )
            legacy = legacy_dir / f"{source.stem}.{suffix}"
            if not legacy.is_file() or legacy.stat().st_size == 0:
                raise RuntimeError(f"Failed to create real .{suffix} smoke fixture")

            output_dir = root / f"converted-{suffix}"
            output_dir.mkdir()
            converted = LibreOfficeHeadlessConverter(executable).convert(
                legacy,
                output_dir,
                target_suffix=f".{suffix}x",
                timeout_seconds=60,
                max_output_bytes=32 * 1024 * 1024,
                max_scratch_bytes=64 * 1024 * 1024,
                max_file_count=256,
            )
            if converted.read_bytes()[:2] != b"PK":
                raise RuntimeError(f"Converted .{suffix}x is not an OOXML package")

        asyncio.run(
            _run_local_docling_chunk_smoke(root, legacy_dir / "legacy-smoke.doc")
        )


async def _run_local_docling_chunk_smoke(root: Path, legacy_doc: Path) -> None:
    """Exercise the real worker preconversion, Docling conversion, and chunk path."""

    legacy_bytes = await asyncio.to_thread(legacy_doc.read_bytes)
    orchestrator = DoclingServeLocalOrchestrator(
        config=LocalOrchestratorConfig(
            num_workers=1,
            shared_models=True,
            scratch_dir=root / "local-worker",
        ),
        converter_manager=build_converter_manager(_build_cm_config()),
    )
    queue_task = asyncio.create_task(orchestrator.process_queue())
    try:
        task = await orchestrator.enqueue(
            sources=[
                DocumentStream(
                    name=legacy_doc.name,
                    stream=BytesIO(legacy_bytes),
                )
            ],
            target=InBodyTarget(),
            task_type=TaskType.CHUNK,
            convert_options=ConvertDocumentsOptions(),
            chunking_options=HierarchicalChunkerOptions(),
            metadata={"tenant_id": "image-smoke"},
        )

        async def wait_until_complete() -> None:
            while not (  # noqa: ASYNC110 - polling the real orchestrator contract
                await orchestrator.task_status(task.task_id)
            ).is_completed():
                await asyncio.sleep(0.1)

        await asyncio.wait_for(wait_until_complete(), timeout=180)
        result = await orchestrator.task_result(task.task_id)
        if result is None or result.num_succeeded != 1:
            raise RuntimeError("Local legacy Office Docling smoke did not succeed")
        chunks = getattr(result.result, "chunks", [])
        if not chunks or not any(
            "legacy office smoke" in chunk.text for chunk in chunks
        ):
            raise RuntimeError(
                "Local legacy Office chunk smoke produced no expected text"
            )
        retained = await orchestrator.task_status(task.task_id)
        if retained.metadata.get("tenant_id") != "image-smoke":
            raise RuntimeError("Local orchestrator lost tenant metadata")
    finally:
        queue_task.cancel()
        await asyncio.gather(queue_task, return_exceptions=True)


if __name__ == "__main__":
    main()
