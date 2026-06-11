"""python-pptx based extraction for the deep-document pipeline.

PPTX flows through the *same* deep-document contract as every other source
(``manifest -> build_deep_document -> S3 object tree``). Docling alone cannot
recover per-shape geometry, run-level styling, speaker notes, or the embedded
image blobs, so for ``.ppt`` / ``.pptx`` we build the normalized manifest
directly from the OOXML via python-pptx instead of from the docling document.

The output shape matches what ``document_builder.normalized_units`` /
``normalized_elements`` expect: a list of ``units`` (one per slide) carrying
``render`` (size + background), ``speakerNotes``, and ``elements`` with ``bbox``
geometry, ``text`` (plain + paragraphs + runs), ``style``, and ``assetRef`` for
images. Image blobs are written next to the deep-document so the existing
``upload_tree`` publisher ships them to S3 with everything else.
"""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

# EMU (English Metric Units) per CSS pixel at 96 DPI. python-pptx geometry is
# all in EMU; the canvas/editor works in pixels, so this is the single
# conversion constant for the whole adapter.
EMU_PER_PX = 9525.0


def manifest_from_pptx(
    pptx_path: str | Path,
    *,
    filename: str | None,
    source_manifest_key: str,
    media_dir: str | Path,
    asset_path_prefix: str,
) -> dict[str, Any]:
    """Build a normalized-extraction manifest for a PPTX file.

    ``media_dir`` is where extracted image blobs are written; ``asset_path_prefix``
    is the path those images are referenced by, relative to the package root
    (so it lines up with ``upload_tree`` S3 keys and ``attach_exported_assets``).
    """
    from pptx import Presentation
    from pptx.util import Emu  # noqa: F401  (imported for type clarity / parity)

    media_root = Path(media_dir)
    presentation = Presentation(str(pptx_path))

    slide_w_emu = int(presentation.slide_width or 0)
    slide_h_emu = int(presentation.slide_height or 0)
    slide_w_px = _emu_to_px(slide_w_emu)
    slide_h_px = _emu_to_px(slide_h_emu)

    assets: list[dict[str, Any]] = []
    asset_index: dict[str, dict[str, Any]] = {}
    units: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides):
        elements: list[dict[str, Any]] = []
        counter = _Counter()
        master_styles = _master_text_styles(slide)
        _walk_shapes(
            shapes=slide.shapes,
            slide_index=slide_index,
            transform=_identity_transform,
            elements=elements,
            counter=counter,
            media_root=media_root,
            asset_path_prefix=asset_path_prefix,
            assets=assets,
            asset_index=asset_index,
            master_styles=master_styles,
        )

        units.append(
            {
                "unitId": f"slide-{slide_index + 1:04d}",
                "unitType": "slide",
                "slideNumber": slide_index + 1,
                "index": slide_index,
                "title": _slide_title(slide) or f"Slide {slide_index + 1}",
                "sourceRefs": {
                    "slideNumber": slide_index + 1,
                    "sourceManifestKey": source_manifest_key,
                },
                "render": {
                    "size": {"px": {"width": slide_w_px, "height": slide_h_px}},
                    "background": _slide_background(slide),
                },
                "speakerNotes": _speaker_notes(slide),
                "elements": elements,
            }
        )

    return {
        "artifactKind": "captify.doclingDocument.normalizedExtraction.v1",
        "source": {
            "originalFileName": filename,
            "fileKind": "presentation",
            "sha256": "",
            "rendererUsed": False,
            "conversionUsed": False,
            "watermarkRisk": False,
            "sourceManifestKey": source_manifest_key,
            "extractor": "python-pptx",
        },
        "units": units,
        "assets": assets,
    }


class _Counter:
    """Monotonic element counter shared across (possibly nested) shape walks."""

    def __init__(self) -> None:
        self.value = 0

    def next(self) -> int:
        self.value += 1
        return self.value


# A transform maps a child-coordinate EMU point/size into absolute slide EMU.
# Represented as (offset_x, offset_y, scale_x, scale_y, child_off_x, child_off_y).
_identity_transform = (0.0, 0.0, 1.0, 1.0, 0.0, 0.0)


def _apply_transform(
    transform: tuple[float, float, float, float, float, float],
    left: float,
    top: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    off_x, off_y, scale_x, scale_y, child_off_x, child_off_y = transform
    abs_left = off_x + (left - child_off_x) * scale_x
    abs_top = off_y + (top - child_off_y) * scale_y
    return abs_left, abs_top, width * scale_x, height * scale_y


def _compose_group_transform(
    parent: tuple[float, float, float, float, float, float],
    group_shape: Any,
) -> tuple[float, float, float, float, float, float]:
    """Derive the child->slide transform for a group, composed with the parent."""
    try:
        xfrm = group_shape._element.grpSpPr.xfrm  # type: ignore[attr-defined]
    except AttributeError:
        xfrm = None
    if xfrm is None:
        return parent

    try:
        off_x = float(xfrm.off.x)
        off_y = float(xfrm.off.y)
        ext_cx = float(xfrm.ext.cx)
        ext_cy = float(xfrm.ext.cy)
        ch_off_x = float(xfrm.chOff.x)
        ch_off_y = float(xfrm.chOff.y)
        ch_ext_cx = float(xfrm.chExt.cx)
        ch_ext_cy = float(xfrm.chExt.cy)
    except (AttributeError, TypeError):
        return parent

    scale_x = (ext_cx / ch_ext_cx) if ch_ext_cx else 1.0
    scale_y = (ext_cy / ch_ext_cy) if ch_ext_cy else 1.0

    # Map the group's own box through the parent transform, then express the
    # child space relative to that absolute box.
    abs_off_x, abs_off_y, _, _ = _apply_transform(parent, off_x, off_y, 0, 0)
    p_scale_x = parent[2]
    p_scale_y = parent[3]
    return (
        abs_off_x,
        abs_off_y,
        scale_x * p_scale_x,
        scale_y * p_scale_y,
        ch_off_x,
        ch_off_y,
    )


def _walk_shapes(
    *,
    shapes: Any,
    slide_index: int,
    transform: tuple[float, float, float, float, float, float],
    elements: list[dict[str, Any]],
    counter: _Counter,
    media_root: Path,
    asset_path_prefix: str,
    assets: list[dict[str, Any]],
    asset_index: dict[str, dict[str, Any]],
    master_styles: dict[str, dict[int, float]],
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            child_transform = _compose_group_transform(transform, shape)
            _walk_shapes(
                shapes=shape.shapes,
                slide_index=slide_index,
                transform=child_transform,
                elements=elements,
                counter=counter,
                media_root=media_root,
                asset_path_prefix=asset_path_prefix,
                assets=assets,
                asset_index=asset_index,
                master_styles=master_styles,
            )
            continue

        element = _shape_element(
            shape=shape,
            shape_type=shape_type,
            slide_index=slide_index,
            transform=transform,
            counter=counter,
            media_root=media_root,
            asset_path_prefix=asset_path_prefix,
            assets=assets,
            asset_index=asset_index,
            master_styles=master_styles,
        )
        if element is not None:
            elements.append(element)


def _shape_element(
    *,
    shape: Any,
    shape_type: Any,
    slide_index: int,
    transform: tuple[float, float, float, float, float, float],
    counter: _Counter,
    media_root: Path,
    asset_path_prefix: str,
    assets: list[dict[str, Any]],
    asset_index: dict[str, dict[str, Any]],
    master_styles: dict[str, dict[int, float]],
) -> dict[str, Any] | None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    bbox = _shape_bbox(shape, transform)
    element_id = f"slide-{slide_index + 1:04d}-element-{counter.next():04d}"
    z_index = counter.value

    is_picture = shape_type in {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }
    if is_picture:
        asset_ref = _extract_image(
            shape=shape,
            media_root=media_root,
            asset_path_prefix=asset_path_prefix,
            assets=assets,
            asset_index=asset_index,
        )
        return {
            "elementId": element_id,
            "type": "image",
            "kind": "picture",
            "bbox": bbox,
            "zIndex": z_index,
            "text": {"plain": "", "paragraphs": [], "runs": []},
            "style": {},
            "assetRef": asset_ref,
            "sourceRefs": {"shapeId": _shape_id(shape)},
        }

    if getattr(shape, "has_table", False):
        return {
            "elementId": element_id,
            "type": "table",
            "kind": "table",
            "bbox": bbox,
            "zIndex": z_index,
            "text": _table_text(shape),
            "style": {},
            "sourceRefs": {"shapeId": _shape_id(shape)},
        }

    if getattr(shape, "has_text_frame", False):
        role = _placeholder_role(shape)
        text = _text_payload(shape.text_frame, role=role, master_styles=master_styles)
        if not text["plain"].strip():
            return None
        return {
            "elementId": element_id,
            "type": "text",
            "kind": "title" if role in {"title", "centerTitle"} else "text",
            "role": role,
            "bbox": bbox,
            "zIndex": z_index,
            "text": text,
            "style": _shape_fill_style(shape),
            "sourceRefs": {"shapeId": _shape_id(shape)},
        }

    # Autoshapes / placeholders with no text and no media are skipped; they add
    # no editable content and would clutter the canvas.
    return None


def _shape_bbox(
    shape: Any, transform: tuple[float, float, float, float, float, float]
) -> dict[str, float]:
    try:
        left = float(shape.left) if shape.left is not None else 0.0
        top = float(shape.top) if shape.top is not None else 0.0
        width = float(shape.width) if shape.width is not None else 0.0
        height = float(shape.height) if shape.height is not None else 0.0
    except (AttributeError, TypeError):
        return {"x": 0.0, "y": 0.0, "w": 0.0, "h": 0.0}

    abs_left, abs_top, abs_w, abs_h = _apply_transform(
        transform, left, top, width, height
    )
    return {
        "x": round(_emu_to_px(abs_left), 2),
        "y": round(_emu_to_px(abs_top), 2),
        "w": round(_emu_to_px(abs_w), 2),
        "h": round(_emu_to_px(abs_h), 2),
    }


def _text_payload(
    text_frame: Any,
    *,
    role: str = "other",
    master_styles: dict[str, dict[int, float]] | None = None,
) -> dict[str, Any]:
    paragraphs: list[dict[str, Any]] = []
    plain_parts: list[str] = []
    flat_runs: list[dict[str, Any]] = []
    styles = master_styles or {}

    for paragraph in text_frame.paragraphs:
        level = int(paragraph.level or 0)
        # The default point size for this paragraph's level, inherited from the
        # slide master's title/body/other text style when a run carries no
        # explicit size (the common case for placeholder text).
        inherited = _inherited_size_pt(styles, role, level)
        runs: list[dict[str, Any]] = []
        for run in paragraph.runs:
            run_payload = _run_payload(run, inherited)
            runs.append(run_payload)
            flat_runs.append(run_payload)
        para_text = "".join(run["text"] for run in runs)
        plain_parts.append(para_text)
        first_size = runs[0]["sizePt"] if runs else inherited
        paragraphs.append(
            {
                "text": para_text,
                "align": _alignment_name(paragraph.alignment),
                "level": level,
                "sizePt": first_size if first_size is not None else inherited,
                "bullet": _paragraph_bullet(paragraph, role, level),
                "runs": runs,
            }
        )

    return {
        "plain": "\n".join(plain_parts).strip(),
        "paragraphs": paragraphs,
        "runs": flat_runs,
    }


def _run_payload(run: Any, inherited_size_pt: float | None = None) -> dict[str, Any]:
    font = run.font
    explicit = _font_size_pt(font)
    return {
        "text": run.text or "",
        "bold": _tri_state(font.bold),
        "italic": _tri_state(font.italic),
        "underline": _tri_state(font.underline),
        # `sizePt` is the *effective* size the renderer should use: the run's own
        # size when set, otherwise the size inherited from the slide master.
        "sizePt": explicit if explicit is not None else inherited_size_pt,
        "explicitSizePt": explicit,
        "color": _font_color(font),
        "font": font.name,
    }


# PowerPoint default text-style sizes (points) by level, used when the master
# carries no explicit `defRPr@sz` for a level. Index 0 == outline level 1.
_DEFAULT_TITLE_SIZES = [44.0]
_DEFAULT_BODY_SIZES = [28.0, 24.0, 20.0, 18.0, 18.0, 18.0, 18.0, 18.0, 18.0]
_DEFAULT_OTHER_SIZES = [18.0]


def _inherited_size_pt(
    master_styles: dict[str, dict[int, float]], role: str, level: int
) -> float | None:
    style_key = (
        "title"
        if role in {"title", "centerTitle"}
        else "body"
        if role in {"body", "subtitle"}
        else "other"
    )
    sizes = master_styles.get(style_key) or {}
    if level in sizes:
        return sizes[level]
    if 0 in sizes:
        return sizes[0]
    # Fall back to PowerPoint's built-in defaults so headers stay headers.
    table = (
        _DEFAULT_TITLE_SIZES
        if style_key == "title"
        else _DEFAULT_BODY_SIZES
        if style_key == "body"
        else _DEFAULT_OTHER_SIZES
    )
    return table[min(level, len(table) - 1)]


def _placeholder_role(shape: Any) -> str:
    """Map a shape to a text role: title / centerTitle / subtitle / body / other."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        if not shape.is_placeholder:
            return "other"
        ph_type = shape.placeholder_format.type
    except (AttributeError, TypeError):
        return "other"
    if ph_type == PP_PLACEHOLDER.TITLE:
        return "title"
    if ph_type == PP_PLACEHOLDER.CENTER_TITLE:
        return "centerTitle"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    return "body"


def _master_text_styles(slide: Any) -> dict[str, dict[int, float]]:
    """Default point sizes per level for the slide master's title/body/other styles."""
    out: dict[str, dict[int, float]] = {}
    try:
        from pptx.oxml.ns import qn

        master = slide.slide_layout.slide_master
        tx_styles = master.element.find(qn("p:txStyles"))
        if tx_styles is None:
            return out
    except (AttributeError, TypeError, Exception):
        return out

    for key, tag in (("title", "p:titleStyle"), ("body", "p:bodyStyle"), ("other", "p:otherStyle")):
        try:
            node = tx_styles.find(qn(tag))
            if node is None:
                continue
            sizes: dict[int, float] = {}
            for level in range(9):
                lvl_pr = node.find(qn(f"a:lvl{level + 1}pPr"))
                if lvl_pr is None:
                    continue
                def_rpr = lvl_pr.find(qn("a:defRPr"))
                if def_rpr is None:
                    continue
                sz = def_rpr.get("sz")
                if sz:
                    sizes[level] = round(int(sz) / 100.0, 2)
            if sizes:
                out[key] = sizes
        except (AttributeError, TypeError, ValueError, Exception):
            continue
    return out


def _paragraph_bullet(paragraph: Any, role: str, level: int) -> dict[str, Any]:
    """Resolve a paragraph's bullet: explicit buChar/buAutoNum/buNone, else by role."""
    try:
        from pptx.oxml.ns import qn

        p_pr = paragraph._p.find(qn("a:pPr"))
    except (AttributeError, TypeError):
        p_pr = None

    if p_pr is not None:
        if p_pr.find(qn("a:buNone")) is not None:
            return {"type": "none"}
        bu_char = p_pr.find(qn("a:buChar"))
        if bu_char is not None:
            return {"type": "bullet", "char": bu_char.get("char") or "\u2022"}
        bu_auto = p_pr.find(qn("a:buAutoNum"))
        if bu_auto is not None:
            return {"type": "number", "scheme": bu_auto.get("type") or "arabicPeriod"}

    # No explicit bullet markup — body placeholders bullet every level by default;
    # titles and standalone text boxes do not. Centered paragraphs are almost
    # always headings, so they stay unbulleted.
    align = _alignment_name(paragraph.alignment)
    if role in {"body", "subtitle"} and align != "center":
        return {"type": "bullet", "char": "\u2022"}
    return {"type": "none"}


def _font_size_pt(font: Any) -> float | None:
    try:
        size = font.size
    except (AttributeError, TypeError):
        return None
    if size is None:
        return None
    try:
        return round(float(size.pt), 2)
    except (AttributeError, TypeError):
        return None


def _font_color(font: Any) -> str | None:
    try:
        color = font.color
        if color is None or color.type is None:
            return None
        rgb = color.rgb
    except (AttributeError, TypeError, Exception):
        return None
    if rgb is None:
        return None
    return f"#{rgb!s}"


def _shape_fill_style(shape: Any) -> dict[str, Any]:
    style: dict[str, Any] = {}
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # MSO_FILL.SOLID
            rgb = fill.fore_color.rgb
            if rgb is not None:
                style["fill"] = f"#{rgb!s}"
    except (AttributeError, TypeError, Exception):
        pass
    return style


def _slide_background(slide: Any) -> dict[str, Any]:
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.type == 1:
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return {"color": f"#{rgb!s}"}
    except (AttributeError, TypeError, Exception):
        pass
    return {"color": "#FFFFFF"}


def _speaker_notes(slide: Any) -> dict[str, str]:
    try:
        if not slide.has_notes_slide:
            return {"raw": "", "cleaned": ""}
        text = slide.notes_slide.notes_text_frame.text or ""
    except (AttributeError, TypeError):
        return {"raw": "", "cleaned": ""}
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return {"raw": text, "cleaned": cleaned}


def _slide_title(slide: Any) -> str | None:
    try:
        title = slide.shapes.title
    except (AttributeError, ValueError):
        title = None
    if title is not None and title.has_text_frame:
        text = title.text_frame.text.strip()
        if text:
            return text.splitlines()[0].strip()
    return None


def _table_text(shape: Any) -> dict[str, Any]:
    rows: list[list[str]] = []
    plain_lines: list[str] = []
    try:
        table = shape.table
    except (AttributeError, TypeError):
        return {"plain": "", "paragraphs": [], "runs": [], "rows": []}
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
        plain_lines.append("\t".join(cells))
    return {
        "plain": "\n".join(plain_lines).strip(),
        "paragraphs": [],
        "runs": [],
        "rows": rows,
    }


def _extract_image(
    *,
    shape: Any,
    media_root: Path,
    asset_path_prefix: str,
    assets: list[dict[str, Any]],
    asset_index: dict[str, dict[str, Any]],
) -> str | None:
    try:
        image = shape.image
        blob = image.blob
        ext = (image.ext or "png").lstrip(".").lower()
    except (AttributeError, TypeError, Exception):
        return None

    digest = hashlib.sha1(blob).hexdigest()[:16]
    filename = f"{digest}.{ext}"
    asset_ref = f"{asset_path_prefix.rstrip('/')}/{filename}"

    if digest not in asset_index:
        media_root.mkdir(parents=True, exist_ok=True)
        target = media_root / filename
        if not target.exists():
            target.write_bytes(blob)
        content_type = image.content_type or f"image/{ext}"
        asset = {
            "assetId": f"asset-{digest}",
            "kind": "embedded_image",
            "role": "display",
            "path": asset_ref,
            "contentType": content_type,
            "sizeBytes": len(blob),
            "display": True,
        }
        assets.append(asset)
        asset_index[digest] = asset

    return asset_ref


def _shape_id(shape: Any) -> str | None:
    try:
        return str(shape.shape_id)
    except (AttributeError, TypeError):
        return None


def _alignment_name(alignment: Any) -> str | None:
    if alignment is None:
        return None
    try:
        return str(alignment).split(".")[-1].split(" ")[0].lower()
    except (AttributeError, TypeError):
        return None


def _tri_state(value: Any) -> bool:
    return bool(value)


def _emu_to_px(value: float) -> float:
    return value / EMU_PER_PX
